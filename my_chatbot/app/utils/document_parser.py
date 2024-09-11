import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def parse_pdf(file_path):
    """
    Parse a PDF document and extract text.
    
    Args:
        file_path (str): Path to the PDF file.
    
    Returns:
        str: Extracted text from the PDF.
    """
    text = ""
    with open(file_path, 'rb') as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def parse_docx(file_path):
    """
    Parse a DOCX document and extract text.
    
    Args:
        file_path (str): Path to the DOCX file.
    
    Returns:
        str: Extracted text from the DOCX.
    """
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file_path):
    """
    Parse a PPTX presentation and extract text.
    
    Args:
        file_path (str): Path to the PPTX file.
    
    Returns:
        str: Extracted text from the PPTX.
    """
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_documents(directory_path):
    """
    Parse all documents in a directory and return their extracted text.
    
    Args:
        directory_path (str): Path to the directory containing documents.
    
    Returns:
        dict: A dictionary with file names as keys and extracted text as values.
    """
    documents = {}
    for filename in os.listdir(directory_path):
        file_path = os.path.join(directory_path, filename)
        if filename.endswith('.pdf'):
            documents[filename] = parse_pdf(file_path)
        elif filename.endswith('.docx'):
            documents[filename] = parse_docx(file_path)
        elif filename.endswith('.pptx'):
            documents[filename] = parse_pptx(file_path)
    return documents